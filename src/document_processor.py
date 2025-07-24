"""
===============================================================================
ESSENTRA - Agentic RAG Chatbot
===============================================================================

Author: Tirumala Manav
Email: tirumalamanav@example.com
GitHub: https://github.com/TirumalaManav
LinkedIn: https://linkedin.com/in/tirumalamanav

Project: ESSENTRA - Advanced Agentic RAG Chatbot
Repository: https://github.com/TirumalaManav/essentra-ai
Created: 2025-07-23
Last Modified: 2025-07-23 17:57:58

License: MIT License
Copyright (c) 2025 Tirumala Manav
"""



import os
import asyncio
import logging
from typing import List, Dict, Any, Optional, Tuple, Union
from pathlib import Path
from datetime import datetime
import uuid
import hashlib

# Core dependencies
import chromadb
from sentence_transformers import SentenceTransformer
import pandas as pd

# Document parsers
import PyPDF2
from docx import Document
from pptx import Presentation

# Async and typing
from pydantic import BaseModel, Field
from concurrent.futures import ThreadPoolExecutor
import aiofiles

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentMetadata(BaseModel):
    """Document metadata model for tracking"""
    file_id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    filename: str
    file_type: str
    file_size: int
    session_id: str
    upload_timestamp: datetime = Field(default_factory=datetime.now)
    processing_status: str = "pending"  # pending, processing, completed, failed
    chunk_count: int = 0
    embedding_model: str = "all-MiniLM-L6-v2"

class DocumentChunk(BaseModel):
    """Individual document chunk model"""
    chunk_id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    content: str
    chunk_index: int
    source_file: str
    metadata: Dict[str, Any] = Field(default_factory=dict)

class ProcessingResult(BaseModel):
    """Result of document processing"""
    success: bool
    file_id: str
    chunks_created: int
    processing_time: float
    error_message: Optional[str] = None
    metadata: DocumentMetadata

class UniversalDocumentProcessor:
    """
    Advanced document processor with session isolation and async processing
    Optimized for production use with comprehensive error handling
    """

    def __init__(self, vector_db_path: str = "./data/vector_db"):
        self.vector_db_path = Path(vector_db_path)
        self.vector_db_path.mkdir(parents=True, exist_ok=True)

        # Initialize ChromaDB client
        self.chroma_client = chromadb.PersistentClient(path=str(self.vector_db_path))

        # Initialize embedding model (async loading)
        self.embedder = None
        self._embedding_model_name = "sentence-transformers/all-MiniLM-L6-v2"

        # Thread pool for CPU-intensive operations
        self.executor = ThreadPoolExecutor(max_workers=4)

        # Supported file types
        self.supported_formats = {
            '.pdf': 'pdf',
            '.docx': 'docx',
            '.doc': 'docx',
            '.pptx': 'pptx',
            '.ppt': 'pptx',
            '.csv': 'csv',
            '.xlsx': 'excel',
            '.xls': 'excel',
            '.txt': 'text',
            '.md': 'markdown'
        }

        logger.info("UniversalDocumentProcessor initialized")

    async def _lazy_load_embedder(self):
        """Lazy load embedding model to improve startup time"""
        if self.embedder is None:
            logger.info(f"Loading embedding model: {self._embedding_model_name}")
            loop = asyncio.get_event_loop()
            self.embedder = await loop.run_in_executor(
                self.executor,
                lambda: SentenceTransformer(self._embedding_model_name)
            )
            logger.info("Embedding model loaded successfully")

    def detect_file_type(self, file_path: Union[str, Path]) -> str:
        """Auto-detect file type from extension with validation"""
        file_path = Path(file_path)
        ext = file_path.suffix.lower()

        if ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {ext}. Supported: {list(self.supported_formats.keys())}")

        return self.supported_formats[ext]

    async def process_document(
        self,
        file_path: Union[str, Path],
        session_id: str,
        chunk_size: int = 500,
        chunk_overlap: int = 50
    ) -> ProcessingResult:
        """
        Main document processing pipeline with async optimization
        """
        start_time = asyncio.get_event_loop().time()
        file_path = Path(file_path)

        try:
            # Validate file exists
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")

            # Create metadata
            file_type = self.detect_file_type(file_path)
            metadata = DocumentMetadata(
                filename=file_path.name,
                file_type=file_type,
                file_size=file_path.stat().st_size,
                session_id=session_id,
                processing_status="processing"
            )

            logger.info(f"Processing document: {file_path.name} (Type: {file_type}, Session: {session_id})")

            # Parse document content
            content = await self._parse_document_async(file_path, file_type)

            # Create chunks
            chunks = self._create_chunks(content, chunk_size, chunk_overlap, metadata.filename)

            # Store in vector database
            await self._store_in_vector_db(chunks, session_id, metadata)

            # Update metadata
            metadata.processing_status = "completed"
            metadata.chunk_count = len(chunks)

            processing_time = asyncio.get_event_loop().time() - start_time

            logger.info(f"Successfully processed {file_path.name}: {len(chunks)} chunks in {processing_time:.2f}s")

            return ProcessingResult(
                success=True,
                file_id=metadata.file_id,
                chunks_created=len(chunks),
                processing_time=processing_time,
                metadata=metadata
            )

        except Exception as e:
            processing_time = asyncio.get_event_loop().time() - start_time
            error_msg = f"Error processing {file_path.name}: {str(e)}"
            logger.error(error_msg)

            return ProcessingResult(
                success=False,
                file_id=metadata.file_id if 'metadata' in locals() else str(uuid.uuid4()),
                chunks_created=0,
                processing_time=processing_time,
                error_message=error_msg,
                metadata=metadata if 'metadata' in locals() else DocumentMetadata(
                    filename=file_path.name,
                    file_type="unknown",
                    file_size=0,
                    session_id=session_id,
                    processing_status="failed"
                )
            )

    async def _parse_document_async(self, file_path: Path, file_type: str) -> str:
        """Async document parsing with proper error handling"""
        loop = asyncio.get_event_loop()

        try:
            if file_type == 'pdf':
                return await loop.run_in_executor(self.executor, self._parse_pdf, file_path)
            elif file_type == 'docx':
                return await loop.run_in_executor(self.executor, self._parse_docx, file_path)
            elif file_type == 'pptx':
                return await loop.run_in_executor(self.executor, self._parse_pptx, file_path)
            elif file_type in ['csv', 'excel']:
                return await loop.run_in_executor(self.executor, self._parse_tabular, file_path)
            elif file_type in ['text', 'markdown']:
                return await self._parse_text(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")

        except Exception as e:
            raise Exception(f"Failed to parse {file_type} file: {str(e)}")

    def _parse_pdf(self, file_path: Path) -> str:
        """Parse PDF files with robust error handling"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                content = []

                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        text = page.extract_text()
                        if text.strip():
                            content.append(f"[Page {page_num + 1}]\n{text}")
                    except Exception as e:
                        logger.warning(f"Error reading page {page_num + 1}: {e}")
                        continue

                if not content:
                    raise ValueError("No readable text found in PDF")

                return "\n\n".join(content)

        except Exception as e:
            raise Exception(f"PDF parsing error: {str(e)}")

    def _parse_docx(self, file_path: Path) -> str:
        """Parse DOCX files with paragraph preservation"""
        try:
            doc = Document(file_path)
            content = []

            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    content.append(text)

            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        content.append(f"[Table] {row_text}")

            if not content:
                raise ValueError("No readable text found in DOCX")

            return "\n\n".join(content)

        except Exception as e:
            raise Exception(f"DOCX parsing error: {str(e)}")

    def _parse_pptx(self, file_path: Path) -> str:
        """Parse PPTX files with slide structure"""
        try:
            prs = Presentation(file_path)
            content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"[Slide {slide_num}]"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())

                if len(slide_content) > 1:  # More than just the slide header
                    content.append("\n".join(slide_content))

            if not content:
                raise ValueError("No readable text found in PPTX")

            return "\n\n".join(content)

        except Exception as e:
            raise Exception(f"PPTX parsing error: {str(e)}")

    def _parse_tabular(self, file_path: Path) -> str:
        """Parse CSV/Excel files with intelligent formatting"""
        try:
            # Determine file type and read accordingly
            if file_path.suffix.lower() == '.csv':
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)

            if df.empty:
                raise ValueError("Empty tabular file")

            # Format as readable text
            content = [f"Data from {file_path.name}"]
            content.append(f"Columns: {', '.join(df.columns.tolist())}")
            content.append(f"Total rows: {len(df)}")
            content.append("")

            # Add sample data (first 10 rows)
            content.append("Sample data:")
            content.append(df.head(10).to_string(index=False))

            # Add summary statistics for numeric columns
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                content.append("\nSummary statistics:")
                content.append(df[numeric_cols].describe().to_string())

            return "\n".join(content)

        except Exception as e:
            raise Exception(f"Tabular file parsing error: {str(e)}")

    async def _parse_text(self, file_path: Path) -> str:
        """Parse text/markdown files with encoding detection"""
        try:
            # Try UTF-8 first, then fallback to other encodings
            encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']

            for encoding in encodings:
                try:
                    async with aiofiles.open(file_path, 'r', encoding=encoding) as file:
                        content = await file.read()
                        if content.strip():
                            return content
                except UnicodeDecodeError:
                    continue

            raise ValueError("Could not decode text file with any supported encoding")

        except Exception as e:
            raise Exception(f"Text parsing error: {str(e)}")

    def _create_chunks(
        self,
        content: str,
        chunk_size: int,
        chunk_overlap: int,
        source_file: str
    ) -> List[DocumentChunk]:
        """Create overlapping chunks with metadata preservation"""
        if not content.strip():
            return []

        # Split by sentences first, then by chunks
        sentences = content.replace('\n', ' ').split('. ')

        chunks = []
        current_chunk = ""
        chunk_index = 0

        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue

            # Check if adding this sentence exceeds chunk size
            if len(current_chunk) + len(sentence) > chunk_size and current_chunk:
                # Create chunk
                chunk = DocumentChunk(
                    content=current_chunk.strip(),
                    chunk_index=chunk_index,
                    source_file=source_file,
                    metadata={
                        "length": len(current_chunk),
                        "sentence_count": current_chunk.count('.') + 1
                    }
                )
                chunks.append(chunk)

                # Start new chunk with overlap
                overlap_text = current_chunk[-chunk_overlap:] if chunk_overlap > 0 else ""
                current_chunk = overlap_text + " " + sentence
                chunk_index += 1
            else:
                current_chunk += " " + sentence if current_chunk else sentence

        # Add the last chunk
        if current_chunk.strip():
            chunk = DocumentChunk(
                content=current_chunk.strip(),
                chunk_index=chunk_index,
                source_file=source_file,
                metadata={
                    "length": len(current_chunk),
                    "sentence_count": current_chunk.count('.') + 1
                }
            )
            chunks.append(chunk)

        return chunks

    async def _store_in_vector_db(
        self,
        chunks: List[DocumentChunk],
        session_id: str,
        metadata: DocumentMetadata
    ):
        """Store chunks in session-isolated ChromaDB collection"""
        if not chunks:
            logger.warning("No chunks to store")
            return

        # Ensure embedder is loaded
        await self._lazy_load_embedder()

        # Get or create session-specific collection
        collection_name = f"session_{session_id}"
        collection = self.chroma_client.get_or_create_collection(
            name=collection_name,
            metadata={"session_id": session_id, "created_at": str(datetime.now())}
        )

        # Prepare data for batch insertion
        documents = [chunk.content for chunk in chunks]
        chunk_ids = [f"{session_id}_{chunk.chunk_id}" for chunk in chunks]
        metadatas = [
            {
                "source_file": chunk.source_file,
                "chunk_index": chunk.chunk_index,
                "file_id": metadata.file_id,
                "upload_timestamp": str(metadata.upload_timestamp),
                **chunk.metadata
            }
            for chunk in chunks
        ]

        # Generate embeddings
        loop = asyncio.get_event_loop()
        embeddings = await loop.run_in_executor(
            self.executor,
            lambda: self.embedder.encode(documents).tolist()
        )

        # Store in ChromaDB
        collection.add(
            documents=documents,
            embeddings=embeddings,
            metadatas=metadatas,
            ids=chunk_ids
        )

        logger.info(f"Stored {len(chunks)} chunks in collection {collection_name}")

    async def retrieve_similar_chunks(
        self,
        query: str,
        session_id: str,
        n_results: int = 5
    ) -> List[Dict[str, Any]]:
        """Retrieve similar chunks for a query with hybrid search"""
        try:
            # Ensure embedder is loaded
            await self._lazy_load_embedder()

            collection_name = f"session_{session_id}"

            try:
                collection = self.chroma_client.get_collection(collection_name)
            except ValueError:
                logger.warning(f"No collection found for session {session_id}")
                return []

            # Generate query embedding
            loop = asyncio.get_event_loop()
            query_embedding = await loop.run_in_executor(
                self.executor,
                lambda: self.embedder.encode([query]).tolist()[0]
            )

            # Search similar chunks
            results = collection.query(
                query_embeddings=[query_embedding],
                n_results=min(n_results, collection.count()),
                include=["documents", "metadatas", "distances"]
            )

            # Format results
            formatted_results = []
            for i in range(len(results['documents'][0])):
                formatted_results.append({
                    "content": results['documents'][0][i],
                    "metadata": results['metadatas'][0][i],
                    "similarity_score": 1 - results['distances'][0][i],  # Convert distance to similarity
                })

            return formatted_results

        except Exception as e:
            logger.error(f"Error retrieving chunks: {str(e)}")
            return []

    def get_session_stats(self, session_id: str) -> Dict[str, Any]:
        """Get statistics for a session"""
        try:
            collection_name = f"session_{session_id}"
            collection = self.chroma_client.get_collection(collection_name)

            # Get all documents to analyze
            all_docs = collection.get(include=["metadatas"])

            # Aggregate stats
            files = set()
            total_chunks = len(all_docs['metadatas'])

            for metadata in all_docs['metadatas']:
                files.add(metadata.get('source_file', 'unknown'))

            return {
                "session_id": session_id,
                "total_files": len(files),
                "total_chunks": total_chunks,
                "files": list(files)
            }

        except ValueError:
            return {
                "session_id": session_id,
                "total_files": 0,
                "total_chunks": 0,
                "files": []
            }

    async def cleanup_session(self, session_id: str) -> bool:
        """Clean up session data"""
        try:
            collection_name = f"session_{session_id}"
            self.chroma_client.delete_collection(collection_name)
            logger.info(f"Cleaned up session: {session_id}")
            return True
        except Exception as e:
            logger.error(f"Error cleaning up session {session_id}: {str(e)}")
            return False

# Testing function for individual component testing
async def test_document_processor():
    """Test the document processor with sample files"""
    processor = UniversalDocumentProcessor()
    session_id = "test_session_123"

    # Test with different file types (you'll need to provide actual files)
    test_files = [
        "./sample_data/3.pdf",
        "./sample_data/1.docx",
        "./sample_data/5.csv"
    ]

    for file_path in test_files:
        if os.path.exists(file_path):
            print(f"\nTesting with: {file_path}")
            result = await processor.process_document(file_path, session_id)
            print(f"Success: {result.success}")
            print(f"Chunks created: {result.chunks_created}")
            print(f"Processing time: {result.processing_time:.2f}s")

            if result.success:
                # Test retrieval
                search_results = await processor.retrieve_similar_chunks(
                    "What is this document about?",
                    session_id,
                    n_results=3
                )
                print(f"Retrieved {len(search_results)} similar chunks")

    # Print session stats
    stats = processor.get_session_stats(session_id)
    print(f"\nSession stats: {stats}")

if __name__ == "__main__":
    # Run the test
    asyncio.run(test_document_processor())
